import os
from pptx import Presentation
from wand.image import Image
from wand.display import display
from wand.drawing import Drawing
from pptx.util import Inches

    
def insert_into_ppt(pr1,img_name):   
    slide1_register = pr1.slide_layouts[1]
    slide1=pr1.slides.add_slide(slide1_register)
    title=slide1.shapes.title
    subtitle=slide1.placeholders[1]
    pic = slide1.shapes.add_picture(os.getcwd()+'/images/'+img_name,Inches(0.5), Inches(3),width=Inches(3), height=Inches(4))
    title.text=img_name
    subtitle.text=img_name


def insert_watermark(img_path,logo_path,output_name):
    with Image(filename=logo_path) as img1:
        img1.resize(1000,500)
        with Image(filename=img_path) as img:
            img.composite_channel('all_channels',img1,'dissolve',0,0)
            img.save(filename=os.getcwd()+'/images/'+output_name)
         

path_imgs=os.listdir(os.getcwd()+'/assets')
path_logo=os.listdir(os.getcwd()+'/watermark')
for i in path_imgs:
    x=os.getcwd()+'/assets/'+i
    y=os.getcwd()+'/watermark/'+path_logo[0]
    insert_watermark(x,y,i)

pr1=Presentation()
water_marked=os.listdir(os.getcwd()+'/images')
for i in water_marked:
    insert_into_ppt(pr1,i)
pr1.save(os.getcwd()+'/ppt/'+'wand_ppt.pptx')
